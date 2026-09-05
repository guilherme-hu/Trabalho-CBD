# -*- coding: utf-8 -*-
"""Aplica as correções técnicas ao DOCX/PPTX preservando os arquivos editáveis.

Dependências: python-docx e python-pptx.
"""
from pathlib import Path

from docx import Document
from docx.shared import Cm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
DOCX = BASE / "DAS_RAID_Trabalho.docx"
PPTX = BASE / "DAS_RAID_Slides.pptx"


def set_paragraph(document, startswith, value):
    for paragraph in document.paragraphs:
        if paragraph.text.strip().startswith(startswith):
            paragraph.text = value
            return paragraph
    for paragraph in document.paragraphs:
        if paragraph.text == value:
            return paragraph
    raise ValueError(f"Parágrafo não encontrado: {startswith}")


def insert_before(anchor, text, style=None):
    paragraph = anchor.insert_paragraph_before(text, style=style)
    return paragraph


def remove_paragraph(paragraph):
    element = paragraph._element
    element.getparent().remove(element)


def update_docx():
    doc = Document(DOCX)
    for paragraph in list(doc.paragraphs):
        if paragraph.text.strip().startswith("[Grupo]:"):
            remove_paragraph(paragraph)
    set_paragraph(doc, "NAS (Network Attached Storage): NAS", "NAS (Network Attached Storage): dispositivo conectado à rede Ethernet/IP e acessado por múltiplos hosts em nível de arquivo, por protocolos como NFS ou SMB/CIFS. O cliente enxerga arquivos e diretórios, não blocos crus.")
    set_paragraph(doc, "SAN (Storage Area Network): SAN", "SAN (Storage Area Network): rede/fabric que apresenta armazenamento compartilhado em nível de bloco, normalmente por Fibre Channel, iSCSI ou NVMe-oF. Compartilhar escrita entre hosts exige coordenação do SGBD ou sistema de arquivos de cluster.")
    set_paragraph(doc, "O SAS (Serial Attached SCSI)", "O SAS (Serial Attached SCSI) é uma interface serial empresarial que transporta o conjunto de comandos SCSI. Continua relevante em HDDs corporativos, enclosures e backplanes tri-mode, com dual-port e expansores; não se presume predominância em novas arquiteturas, nas quais NVMe e armazenamento distribuído também são comuns.")
    set_paragraph(doc, "O USB 3.0", "O USB 3.0/USB 3.2 evoluiu de 5 para 10 e 20 Gbps. O USB4 Version 2.0 prevê enlaces de até 80 Gbps; o modo de 120 Gbps é opcional e assimétrico. Host, cabo e dispositivo precisam negociar a mesma capacidade, e taxa nominal não é vazão sustentada de armazenamento.")
    set_paragraph(doc, "Diferentemente das interfaces anteriores, o HBA", "HBA (Host Bus Adapter) é uma categoria de adaptador/controladora, não um protocolo universal: liga o host, geralmente por PCIe, a SAS, Fibre Channel ou outro transporte. USB usa host controller e SATA pode estar integrado ao chipset; nem toda E/S passa por uma placa HBA dedicada.")
    set_paragraph(doc, "A escolha do tamanho do stripe", "A escolha da stripe unit (chunk) deve ser alinhada ao tamanho de página/I/O, ao número de discos e ao workload. Stripe é o conjunto de unidades na mesma posição lógica; stripe width é a quantidade de dados por stripe antes da paridade. Unidades menores não favorecem OLTP automaticamente: podem espalhar uma única E/S por vários discos e aumentar trabalho de paridade.")
    set_paragraph(doc, 'Shadowing ("espelhamento em sombra")', "Disk shadowing foi usado historicamente para espelhamento físico, inclusive em sistemas Tandem, mas o termo não é sinônimo universal de RAID 1. Shadow paging, snapshots e cópias lógicas do SGBD possuem semânticas distintas e não devem ser classificados como mirroring de disco.")
    set_paragraph(doc, "Essa fórmula evidencia três relações", "A aproximação evidencia tendências, mas só sob discos idênticos, falhas independentes e exponenciais, detecção imediata, reparo perfeito e unidades consistentes. Ela não inclui URE, setores latentes, falhas correlacionadas nem carga de rebuild. Disponibilidade é outra métrica: A = MTTF/(MTTF+MTTR). Exemplo: MTTF de 1.000.000 h e MTTR de 24 h resulta em A ≈ 99,9976%, sem transformar o valor em garantia.")
    set_paragraph(doc, "Backup e arquivamento de longo prazo", "Backup e arquivamento de longo prazo: RAID não é backup. Use cópia independente, versionada e preferencialmente off-site/imutável, com RPO, RTO e restauração testada. RAID 6 pode aumentar disponibilidade do repositório, mas não protege contra exclusão, corrupção propagada, ransomware ou perda do site/controlador.")
    set_paragraph(doc, "Quanto à interface e à controladora", "Quanto à interface/controladora, a escolha deve seguir mídia, fila, caminhos, flush/FUA, suporte e workload. SAS é uma opção empresarial; NVMe e software-defined storage também podem ser adequados. Cache write-back só é aceitável se protegido por bateria/flash e se respeitar flushes. USB4 é apropriado sobretudo a transporte, desenvolvimento e backup externo, não automaticamente a dados primários de produção.")

    interfaces = [
        ["Interface", "Taxa e latência", "Distância / hot-plug", "Duplex, filas e dispositivos", "Confiabilidade / estado", "Uso em SBD"],
        ["SATA III", "6 Gbps; até ~600 MB/s; latência da mídia", "~1 m; hot-plug depende de AHCI/backplane", "Ponto a ponto; NCQ até 32", "Atual; proteção fim a fim limitada", "Capacidade, entrada, dev"],
        ["eSATA", "3/6 Gbps", "Até ~2 m; hot-plug dependente", "Ponto a ponto", "Legado; sem energia no cabo", "Disco externo legado"],
        ["SAS-4", "22,5 Gbps; ~2.400 MB/s por lane", "Até ~10 m; hot-plug", "Full-duplex, dual-port, expansores, filas profundas", "Atual em HDD/backplane empresarial", "Arrays e caminhos redundantes"],
        ["FireWire 800", "800 Mbps; ciclo isócrono não é latência de storage", "~4,5 m cobre; hot-plug", "Peer-to-peer; até 63 nós", "Descontinuado/legado", "Nenhum projeto novo"],
        ["USB 3.2 Gen2x2", "20 Gbps nominal; latência variável", "~1 m; hot-plug", "Host-cêntrico; até 127; UASP quando suportado", "Ponte/driver podem afetar flush/FUA", "Backup/transporte/dev"],
        ["USB4 v2", "80 Gbps; 120 opcional e assimétrico", "Cabo certificado; hot-plug", "Túneis dinâmicos; capacidade mútua", "Atual, dependente do ecossistema", "SSD externo compatível"],
        ["HBA", "Depende do protocolo e PCIe", "Dentro do host", "SAS/FC etc.; não é protocolo", "RAID exige cache protegido e flush correto", "Conectar SAS/FC; modo IT ou RAID"],
    ]
    table = doc.tables[1]
    for r, row in enumerate(interfaces):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    raid = [
        ["Nível", "Mínimo / útil", "Dados e paridade", "Falhas e rebuild", "Leitura / escrita", "Vantagens e limitações", "Uso atual"],
        ["0", "2 / N×", "Striping; sem paridade", "0; qualquer falha perde volume", "Altas / altas", "Capacidade e desempenho; sem proteção", "Scratch recriável apenas"],
        ["1", "2 / 50%", "Espelho", "1 por par; copia do sobrevivente", "Boa / 2 gravações", "Simples; custo de 50%", "Logs/sistema com backup"],
        ["2", "Vários / variante", "Bits + código Hamming dedicado", "Depende do código", "Lockstep", "Histórico; discos já fazem ECC", "Obsoleto"],
        ["3", "3 / (N−1)×", "Byte/fino + paridade dedicada", "1; todos participam", "Sequencial boa / gargalo", "Bom stream; baixa concorrência", "Histórico"],
        ["4", "3 / (N−1)×", "Bloco + paridade dedicada", "1; paridade é gargalo", "Leituras paralelas / escrita penalizada", "Simples; gargalo central", "Nicho/histórico"],
        ["5", "3 / (N−1)×", "Bloco + paridade distribuída", "1; janela longa e URE importam", "Boa / RMW 4 I/Os no caso pequeno", "Eficiência; risco em discos grandes", "Somente após análise de rebuild"],
        ["6", "4 / (N−2)×", "Bloco + dupla paridade", "2; rebuild ainda degrada", "Boa / RMW até 6 I/Os no caso pequeno", "Mais proteção; escrita/capacidade", "Arrays grandes, condicionado"],
        ["0+1", "4 / 50%", "Espelho de stripes", "Depende da localização; recopia conjunto", "Alta / alta", "Segunda falha no conjunto sobrevivente perde", "RAID 10 preferível"],
        ["10", "4 / 50%", "Stripe de pares espelhados", "Várias, se em pares distintos", "Alta / alta", "Rebuild localizado; custo 50%", "OLTP, após RPO/RTO"],
    ]
    table = doc.tables[2]
    for r, row in enumerate(raid):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    # Registra as correções no Post-Mortem sem inventar autoria.
    errors = doc.tables[6]
    errors.cell(5, 1).text = "RAID foi confundido com backup; MTTDL omitia hipóteses; shadowing foi tratado como sinônimo universal."
    errors.cell(5, 2).text = "Auditoria cruzada do relatório, slides e fontes primárias."
    errors.cell(5, 3).text = "Recomendação corrigida; hipóteses/availability adicionadas; conceitos separados."

    if not any(p.text.startswith("10.1 Política de proteção") for p in doc.paragraphs):
        anchor = next(p for p in doc.paragraphs if p.text.startswith("11. Conclusão"))
        insert_before(anchor, "10.1 Política de proteção e decisão")
        insert_before(anchor, "Toda recomendação deve registrar workload, tamanho/alinhamento de I/O, capacidade, janela de rebuild, cache protegido, RPO e RTO. RAID fornece disponibilidade contra falhas previstas; backup fornece retorno a um estado anterior. Hot spare reduz tempo de exposição, mas não é backup. Logs exigem write-through ou write-back protegido e respeito a flush/FUA. Tempdb em RAID 0 só é aceitável se a instância puder reiniciar e reconstruir o conteúdo dentro do RTO.")

    refs = next(p for p in doc.paragraphs if p.text.startswith("13. Referências"))
    for source in reversed([
        "USB-IF. USB4. https://www.usb.org/usb4. Acesso em 05 set. 2026.",
        "INTEL. Intel RAID Software User Guide. https://cdrdv2-public.intel.com/840787/SWUG_full_featured_entry_RAID.pdf. Acesso em 05 set. 2026.",
        "NETAPP. Default RAID policies for ONTAP local tiers. https://docs.netapp.com/us-en/ontap/disks-aggregates/default-raid-policies-aggregates-concept.html. Acesso em 05 set. 2026.",
        "NETAPP. TR-3298: RAID-DP. https://www.netapp.com/media/19939-tr-3298.pdf. Acesso em 05 set. 2026.",
    ]):
        if not any(source in p.text for p in doc.paragraphs):
            insert_before(refs, source)

    doc.save(DOCX)


def replace_in_pptx(prs, old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    combined = "".join(run.text for run in paragraph.runs)
                    if old in combined and paragraph.runs:
                        paragraph.runs[0].text = combined.replace(old, new)
                        for run in paragraph.runs[1:]:
                            run.text = ""
                        continue
                    for run in paragraph.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def replace_exact_shape_text(slide, old, new):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip() == old:
            paragraph = shape.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = new
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = new


def add_text(slide, text, x, y, w, h, size=18, color="18243A", bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    return box


def add_correction_slide(prs, title, left_title, left, right_title, right, position):
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(246, 248, 251)
    add_text(slide, title, 0.75, 0.45, 11.8, 0.6, 26, "18243A", True)
    add_text(slide, left_title, 0.85, 1.45, 5.6, 0.4, 16, "C58D24", True)
    add_text(slide, left, 0.85, 1.9, 5.6, 4.7, 15)
    add_text(slide, right_title, 6.85, 1.45, 5.6, 0.4, 16, "C58D24", True)
    add_text(slide, right, 6.85, 1.9, 5.6, 4.7, 15)
    add_text(slide, "Fontes e hipóteses completas no relatório corrigido.", 0.85, 6.8, 11.6, 0.3, 10, "667085")
    sld_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(sld_id)
    prs.slides._sldIdLst.insert(position, sld_id)


def update_pptx():
    prs = Presentation(PPTX)
    replacements = {
        "Padrão em SBD": "Relevante em HDD/backplanes",
        "80 Gbps simétricos / 120 Gbps assimétricos": "80 Gbps; 120 Gbps opcional e assimétrico (host/cabo/dispositivo compatíveis)",
        '"Shadowing" é termo histórico (ex.: Tandem NonStop) equivalente a mirroring — mesma técnica, nomes diferentes.': "Disk shadowing pode nomear espelhamento histórico; shadow paging, snapshots e cópias lógicas não são RAID 1.",
        "Uso recomendado: uso geral em SBD — melhor equilíbrio custo × desempenho × confiabilidade": "Uso condicionado: avaliar URE, duração/carga de rebuild, escrita pequena e RPO/RTO",
        "RAID 6\nSATA": "Backup independente/versionado\nRAID 6 apenas para disponibilidade do repositório",
        "Quadro completo (9 níveis) no relatório em PDF, Tabela 2.": "Quadro completo corrigido no relatório, incluindo rebuild, leitura/escrita e relevância atual.",
    }
    for old, new in replacements.items():
        replace_in_pptx(prs, old, new)

    recommendations = prs.slides[15]
    replace_exact_shape_text(recommendations, "RAID 6", "Backup independente/versionado")
    replace_exact_shape_text(recommendations, "SATA", "RAID 6 só dá disponibilidade")
    replace_exact_shape_text(recommendations, "RAID 0", "RAID 0 só se recriável no RTO")

    titles = {shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)}
    if "RAID fornece disponibilidade; backup fornece recuperação" not in titles:
        add_correction_slide(
            prs, "RAID fornece disponibilidade; backup fornece recuperação",
            "RAID protege contra falhas previstas",
            "• Tolera falha de disco conforme o nível\n• Pode manter o serviço durante rebuild\n• Hot spare reduz a janela de exposição\n• Não preserva estado histórico\n• Corrupção/exclusão pode atingir todo o array",
            "Backup é uma cópia independente",
            "• Versionado e preferencialmente off-site/imutável\n• Protege contra erro lógico e ransomware\n• Precisa de RPO e RTO\n• Restauração deve ser testada\n• RAID pode proteger o repositório, mas não substituí-lo",
            16,
        )
        add_correction_slide(
            prs, "Confiabilidade: hipótese antes da fórmula",
            "O que o MTTDL simplificado assume",
            "• Discos idênticos\n• Falhas independentes/exponenciais\n• Detecção imediata\n• Reparo perfeito\n• Unidades consistentes\n\nNão cobre URE, setores latentes, correlação nem carga de rebuild.",
            "Métrica complementar",
            "Disponibilidade: A = MTTF/(MTTF + MTTR)\n\nExemplo: 1.000.000 h e MTTR 24 h → 99,9976%.\n\nValor médio não é garantia. Dimensionar com RPO/RTO e teste de falha.",
            16,
        )
        add_correction_slide(
            prs, "Níveis históricos e diferença 0+1 × 10",
            "RAID 2, 3 e 4",
            "RAID 2: striping em bits + código Hamming; histórico e dependente da variante.\n\nRAID 3: granularidade fina/byte + paridade dedicada; discos em lockstep.\n\nRAID 4: blocos + paridade dedicada; leituras paralelas, gargalo na escrita.",
            "RAID 0+1 e RAID 10",
            "0+1 perde um stripe após a primeira falha; a tolerância adicional depende de onde falha o próximo disco.\n\nRAID 10 expõe somente o par afetado; perde dados apenas se ambos os membros do mesmo par falharem.",
            16,
        )
    prs.save(PPTX)


if __name__ == "__main__":
    update_docx()
    update_pptx()
    print("DOCX e PPTX de DAS/RAID atualizados.")
