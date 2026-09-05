# -*- coding: utf-8 -*-
"""Aplica as correções técnicas ao DOCX/PPTX da proposta AWS.

Dependências: python-docx e python-pptx.
"""
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
DOCX = BASE / "Cloud_Computing_AWS_Trabalho.docx"
PPTX_OLD = BASE / "Cloud_Computing_AWS_Slides(1).pptx"
PPTX = BASE / "Cloud_Computing_AWS_Slides.pptx"


def set_paragraph(document, startswith, value):
    for paragraph in document.paragraphs:
        if paragraph.text.strip().startswith(startswith):
            paragraph.text = value
            return paragraph
    for paragraph in document.paragraphs:
        if paragraph.text == value:
            return paragraph
    raise ValueError(f"Parágrafo não encontrado: {startswith}")


def insert_before(anchor, text, style=None):
    return anchor.insert_paragraph_before(text, style=style)


def remove_paragraph(paragraph):
    element = paragraph._element
    element.getparent().remove(element)


def add_table_before(document, anchor, rows):
    table = document.add_table(rows=len(rows), cols=len(rows[0]))
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    anchor._p.addprevious(table._tbl)
    return table


def table_by_header(document, header):
    for table in document.tables:
        if table.rows and table.cell(0, 0).text.strip() == header:
            return table
    raise ValueError(f"Tabela não encontrada: {header}")


def update_docx():
    doc = Document(DOCX)
    for paragraph in list(doc.paragraphs):
        if paragraph.text.strip().startswith("[Grupo]:"):
            remove_paragraph(paragraph)
    set_paragraph(doc, "Cloud Storage (armazenamento em nuvem)", "Cloud Storage é a família de serviços que persiste dados abstraindo parte da infraestrutura física. Inclui objetos (S3), blocos (EBS), arquivos (EFS), classes de arquivamento (S3 Glacier) e armazenamento gerenciado de bancos/warehouses. O acesso pode ocorrer por Internet, VPN, Direct Connect ou endpoints privados. Semântica, durabilidade, disponibilidade, latência e responsabilidade variam por serviço.")
    set_paragraph(doc, "SaaS (Software as a Service)", "SaaS (Software as a Service) — aplicação completa operada pelo provedor; o cliente ainda administra usuários, conteúdo, classificação e configurações permitidas. Exemplo no ecossistema AWS: Amazon QuickSight.")
    set_paragraph(doc, "DBaaS (Database as a Service)", "DBaaS (Database as a Service) — o provedor opera infraestrutura, engine e tarefas gerenciadas conforme a modalidade; o cliente continua responsável por dados, esquema, consultas, IAM/credenciais, classificação e configurações disponíveis. Exemplo: Amazon RDS.")
    set_paragraph(doc, "DWaaS (Data Warehouse as a Service)", "DWaaS (Data Warehouse as a Service) — serviço analítico gerenciado com armazenamento colunar/MPP. Comparado a um DW autogerenciado em EC2, transfere ao provedor engine e infraestrutura, mas mantém com o cliente dados, modelo, cargas, consultas, acesso e governança. Exemplo: Amazon Redshift.")
    set_paragraph(doc, "Nuvem pública:", "Nuvem pública: infraestrutura operada por provedor para múltiplos clientes, com isolamento lógico e serviços padronizados. A classificação decorre de propriedade/operação/tenancy, não da obrigação de usar Internet pública: VPN, Direct Connect e endpoints privados também podem ser usados.")
    set_paragraph(doc, "Escalabilidade é a capacidade", "Escalabilidade é a capacidade de suportar crescimento vertical ou horizontal. Elasticidade é ajustar recursos rapidamente — idealmente de forma automática — para acompanhar aumento e redução da demanda. Um sistema pode ser escalável sem ser elástico.")
    set_paragraph(doc, "O Amazon EC2", "O Amazon EC2 fornece instâncias virtuais sob demanda. On-Demand cobra pelo uso; Spot usa capacidade ociosa interrompível; Savings Plans oferecem desconto mediante compromisso de uso; Reserved Instances são principalmente benefício de faturamento, e somente modalidades zonais específicas incluem reserva de capacidade. Bancos autogerenciados exigem gestão do SO, engine, backup, HA e patches pelo cliente.")
    set_paragraph(doc, "Amazon EBS", "Amazon EBS (Elastic Block Store): volumes em bloco anexados a EC2, com tipos e limites próprios de capacidade/IOPS/throughput. É adequado a sistemas de arquivos e bancos autogerenciados quando a arquitetura respeita limites por volume/instância e requisitos de disponibilidade.")
    set_paragraph(doc, "Amazon S3 Glacier", "Amazon S3 Glacier: classes de dados frios dentro do S3. Instant Retrieval mantém acesso em milissegundos; Flexible Retrieval oferece opções de minutos a horas; Deep Archive tipicamente exige 12 a 48 horas. Cada classe possui cobrança de recuperação, metadados e permanência mínima.")
    set_paragraph(doc, "Amazon RDS", "Amazon RDS: DBaaS para engines relacionais. Backups e patches são gerenciados conforme configuração; replicação/failover automático dependem de Multi-AZ ou modalidade equivalente explicitamente contratada, não de toda instância RDS.")
    set_paragraph(doc, "Amazon Aurora", "Amazon Aurora: engine compatível com MySQL/PostgreSQL com armazenamento distribuído próprio. Desempenho, disponibilidade, custo e limite de cluster variam por versão, região e configuração; comparar por benchmark e SLA em vez de afirmar superioridade universal.")
    set_paragraph(doc, "Amazon DynamoDB", "Amazon DynamoDB: banco NoSQL chave-valor/documento gerenciado. A AWS descreve latência de milissegundos de um dígito para a maioria das operações unitárias; desenho de chave, tamanho do item, capacidade e rede ainda determinam a latência observada.")
    set_paragraph(doc, "Nesta seção, estima-se", "Nesta seção, estima-se um cenário analítico de 3,5 PB decimais na região us-east-1, com preços de lista consultados em 05/09/2026. Não é cotação. Workload-base: data lake em S3, objetos colunares/particionados, crescimento e frequência de leitura a serem confirmados pelo grupo; descontos negociados, suporte e impostos ficam fora até cotação.")
    set_paragraph(doc, "Premissa de conversão", "Premissa de conversão: PB é unidade decimal. Logo, 3,5 PB = 3.500 TB = 3.500.000 GB para cobrança em GB. Se a intenção fosse 3,5 PiB, seriam 3.670.016 GiB ou aproximadamente 3.940.649,674 GB; as bases não podem ser misturadas.")
    set_paragraph(doc, "O Amazon S3 Standard utiliza", "O Amazon S3 Standard usa faixas em us-east-1: US$0,023/GB nos primeiros 50 TB, US$0,022/GB nos 450 TB seguintes e US$0,021/GB acima de 500 TB. Para 3.500.000 GB: 50.000×0,023 + 450.000×0,022 + 3.000.000×0,021 = US$74.050/mês, armazenamento apenas.")
    set_paragraph(doc, "O custo de armazenamento é apenas", "Armazenamento é apenas uma parcela. O TCO deve incluir PUT/GET/LIST, transições, recuperação, mínimo de 128 KB em Standard-IA, metadados de arquivo, permanência mínima, KMS, catálogo, observabilidade, replicação, suporte e transferência. Egress depende de volume/destino/faixa e, em volumes altos, pode exigir cotação; não se apresenta total sem essas premissas.")
    set_paragraph(doc, "Para processar uma base dessa escala", "Para analytics, um desenho plausível é S3 em Parquet/ORC com partições, Glue Catalog e Athena/EMR/Redshift Spectrum. A US$5/TB escaneado, uma varredura integral de 3.500 TB custa US$17.500; uma por dia por 30 dias, US$525.000/mês; ler 10% após compressão/poda custa US$1.750 por execução. Somar requests e resultados.")
    set_paragraph(doc, "Em síntese, a AWS", "Em síntese, a AWS comporta 3,5 PB em S3, mas não se deve presumir um único Aurora, RDS ou volume EBS: limites por cluster/volume/instância exigem validação por versão/região. Armazenamento puro varia de US$74.050/mês em Standard a US$3.465/mês em Deep Archive sob as tarifas assumidas, porém o TCO depende do padrão de acesso, processamento, recuperação, egress, crescimento, RPO/RTO e governança.")

    responsibility = [
        ["Camada", "IaaS", "PaaS", "SaaS", "DBaaS/DWaaS"],
        ["Aplicação/código/consultas", "Cliente", "Cliente", "Provedor; cliente configura/usa", "Cliente: esquema, cargas e consultas"],
        ["Dados, IAM e governança", "Cliente", "Cliente", "Cliente para seu conteúdo/acesso", "Cliente"],
        ["Runtime/engine, SO e middleware", "Cliente", "Provedor", "Provedor", "Provedor conforme o serviço"],
        ["Virtualização/infra física", "Provedor", "Provedor", "Provedor", "Provedor"],
    ]
    table = table_by_header(doc, "Camada")
    for r, row in enumerate(responsibility):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    storage_cost = [
        ["Faixa", "Volume (TB)", "Preço/GB (US$)", "Custo mensal (US$)"],
        ["Primeiros 50 TB", "50", "0,023", "1.150"],
        ["50–500 TB", "450", "0,022", "9.900"],
        ["Acima de 500 TB", "3.000", "0,021", "63.000"],
        ["Total S3 Standard", "3.500", "—", "74.050/mês; 888.600/ano"],
    ]
    table = table_by_header(doc, "Faixa")
    for r, row in enumerate(storage_cost):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    classes = [
        ["Classe", "US$/GB-mês", "Armazenamento mensal", "Condição principal"],
        ["S3 Standard", "faixas 0,023/0,022/0,021", "US$74.050", "Acesso frequente; requests à parte"],
        ["Standard-IA", "0,0125", "US$43.750", "Retrieval, 128 KB mínimo, 30 dias"],
        ["Glacier Flexible", "0,0036", "US$12.600", "Retrieval/metadados; 90 dias"],
        ["Glacier Deep Archive", "0,00099", "US$3.465", "12–48 h; retrieval/metadados; 180 dias"],
    ]
    table = table_by_header(doc, "Classe")
    for r, row in enumerate(classes):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    # Completa a seção AWS com armazenamento em arquivo.
    if not any(p.text.startswith("Amazon EFS") for p in doc.paragraphs):
        anchor = next(p for p in doc.paragraphs if p.text.startswith("7.3 Banco"))
        insert_before(anchor, "Amazon EFS (Elastic File System): serviço de arquivos NFS elástico e compartilhável entre múltiplos clientes. Difere de S3 (objeto/API) e EBS (bloco anexado a instância); latência, throughput e preço dependem da classe/modo.")

    if not any(p.text.startswith("8.5 Arquitetura") for p in doc.paragraphs):
        anchor = next(p for p in doc.paragraphs if p.text.startswith("9. Conclusão"))
        insert_before(anchor, "8.5 Arquitetura, limites e cenários")
        insert_before(anchor, "Arquitetura-base: S3 como data lake em formato colunar/particionado; Glue Catalog; Athena para ad hoc, EMR para lote e Redshift/Spectrum para analytics recorrente; lifecycle por temperatura; camada separada de banco operacional. S3 não limita o total agregado do bucket, mas cada serviço consumidor tem quotas. Aurora e EBS possuem limites por cluster/volume e não acomodam 3,5 PB em uma única unidade. Confirmar limites atuais por versão e região.")
        add_table_before(doc, anchor, [
            ["Cenário", "Acesso", "Arquitetura", "Custos a modelar"],
            ["Econômico/arquivo", "Raro; RTO 12–48 h", "S3 Glacier + catálogo/imutabilidade", "Storage, metadata, retrieval, permanência"],
            ["Analytics típico", "Consultas particionadas", "S3 Parquet + Glue + Athena/EMR", "Storage, TB lido, requests, suporte"],
            ["Alto desempenho", "Recorrente/baixa latência", "S3 + Redshift/compute dimensionado", "Compute, réplicas, cache, AZ/egress"],
        ])
        insert_before(anchor, "Transferência inicial ideal: 3,5 PB a 10 Gb/s ≈ 32,4 dias; a 100 Gb/s ≈ 3,24 dias, antes de protocolo, criptografia e contenção. Avaliar Direct Connect, DataSync e dispositivos de transferência com quotas/cotações atuais. Para cada cenário, definir crescimento, compressão, tamanho/número de objetos, RPO, RTO e teste de recuperação.")

    # Há várias tabelas iniciadas por “#”; esta é a única com quatro colunas.
    errors = next(
        table for table in doc.tables
        if len(table.columns) == 4
        and table.cell(0, 1).text.strip() == "Risco / erro identificado"
    )
    errors.cell(2, 1).text = "PB decimal foi confundido com PiB/GiB e tabelas usaram bases diferentes."
    errors.cell(2, 2).text = "Reprodução de cada multiplicação e conferência das unidades AWS."
    errors.cell(2, 3).text = "Base única: 3,5 PB = 3.500.000 GB; todas as classes e gráficos recalculados."
    errors.cell(5, 1).text = "Custo de processamento e viabilidade por serviço estavam ausentes."
    errors.cell(5, 2).text = "Cálculo Athena e consulta aos limites oficiais de S3/Aurora/EBS."
    errors.cell(5, 3).text = "Cenários, arquitetura de data lake, limites, RPO/RTO e TCO adicionados."

    # Restaura a declaração individual caso uma versão anterior do gerador
    # tenha escrito por índice na tabela errada após inserir a tabela de cenários.
    review = next(
        table for table in doc.tables
        if len(table.columns) == 4
        and table.cell(0, 1).text.strip().startswith("Declaro que revisei")
    )
    corrupted_review = any(
        "PB decimal foi confundido" in cell.text
        or "Custo de processamento" in cell.text
        for row in review.rows[1:]
        for cell in row.cells
    )
    if corrupted_review:
        for index in range(1, 6):
            review.cell(index, 0).text = f"[Nome {index}]"
            review.cell(index, 1).text = "[ ] Sim"
            review.cell(index, 2).text = "[__/__/2026]"
            review.cell(index, 3).text = "________________"

    refs = next(p for p in doc.paragraphs if p.text.startswith("11. Referências"))
    sources = [
        "NIST. SP 800-145: The NIST Definition of Cloud Computing. https://csrc.nist.gov/pubs/sp/800/145/final. Acesso em 05 set. 2026.",
        "AWS. Amazon S3 Pricing. https://aws.amazon.com/s3/pricing/. Acesso em 05 set. 2026.",
        "AWS. Amazon Athena Pricing. https://aws.amazon.com/athena/pricing/. Acesso em 05 set. 2026.",
        "AWS. What is Amazon EFS? https://docs.aws.amazon.com/efs/latest/ug/whatisefs.html. Acesso em 05 set. 2026.",
        "AWS. Aurora quotas and limits. https://docs.aws.amazon.com/AmazonRDS/latest/AuroraUserGuide/CHAP_Limits.html. Acesso em 05 set. 2026.",
        "AWS. EBS volume constraints. https://docs.aws.amazon.com/ebs/latest/userguide/volume_constraints.html. Acesso em 05 set. 2026.",
        "AWS. Shared Responsibility Model. https://aws.amazon.com/compliance/shared-responsibility-model/. Acesso em 05 set. 2026.",
    ]
    for source in reversed(sources):
        if not any(source in p.text for p in doc.paragraphs):
            insert_before(refs, source)

    doc.save(DOCX)


def replace_in_pptx(prs, old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    combined = "".join(run.text for run in paragraph.runs)
                    if old in combined and paragraph.runs:
                        paragraph.runs[0].text = combined.replace(old, new)
                        for run in paragraph.runs[1:]:
                            run.text = ""
                        continue
                    for run in paragraph.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def replace_exact_paragraph(slide, old, new):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            combined = "".join(run.text for run in paragraph.runs)
            if combined.strip() != old:
                continue
            if paragraph.runs:
                paragraph.runs[0].text = new
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = new


def add_text(slide, text, x, y, w, h, size=18, color="1D275C", bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    return box


def add_slide(prs, title, left_title, left, right_title, right, position):
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 249, 253)
    add_text(slide, title, 0.75, 0.45, 11.8, 0.65, 26, "1D275C", True)
    add_text(slide, left_title, 0.85, 1.45, 5.6, 0.4, 16, "E28A00", True)
    add_text(slide, left, 0.85, 1.9, 5.6, 4.75, 15)
    add_text(slide, right_title, 6.85, 1.45, 5.6, 0.4, 16, "E28A00", True)
    add_text(slide, right, 6.85, 1.9, 5.6, 4.75, 15)
    add_text(slide, "Premissas, fontes e fórmulas completas no relatório.", 0.85, 6.8, 11.6, 0.3, 10, "667085")
    sld_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(sld_id)
    prs.slides._sldIdLst.insert(position, sld_id)


def update_pptx():
    source = PPTX if PPTX.exists() else PPTX_OLD
    prs = Presentation(source)
    replacements = {
        "Subconjunto da nuvem voltado ao armazenamento remoto de dados, replicados entre discos/zonas e acessados via APIs, sem que o cliente precise saber em qual disco físico os dados residem.": "Família de serviços de objeto (S3), bloco (EBS), arquivo (EFS), archive e bancos gerenciados. Pode usar Internet ou conectividade privada; semântica e SLA variam.",
        "Gmail, M365": "Amazon QuickSight",
        "infraestrutura de terceiros, compartilhada (multi-tenant), acesso via Internet.": "infraestrutura operada por provedor/multi-tenant; Internet não é obrigatória (VPN, Direct Connect, endpoints privados).",
        "Vertical (scale up) e horizontal (scale out)": "Escalabilidade = crescer vertical/horizontal; elasticidade = ajustar rapidamente para cima e para baixo",
        "• Reserved Instances — desconto por 1–3 anos": "• Reserved Instances — benefício de faturamento; capacidade só em modalidades zonais específicas",
        "Armazenamento — S3, EBS e Glacier": "Armazenamento — S3, EBS, EFS e Glacier",
        "backup e failover automáticos": "backup gerenciado; failover se Multi-AZ/configurado",
        "≈ 3.584 TB ≈ 3.670.016 GB (padrão binário: 1 PB = 1.024 TB)": "3.500 TB = 3.500.000 GB (PB decimal; não confundir com PiB)",
        "• Preços de lista, ago/2026": "• Preços de lista consultados em 05/09/2026",
        "• Custo de armazenamento apenas": "• Storage apenas; requests/retrieval/egress/compute separados",
        "S3 Standard → Glacier Deep Archive representa economia de armazenamento superior a 95%. Mas Glacier Deep Archive leva 12–48h para recuperar dados — inviável para consultas ativas. Na prática: camadas combinadas (lifecycle policies), não uma classe única.": "Classes frias reduzem storage, mas incluem retrieval, metadados e permanência mínima. Deep Archive: 12–48 h. Escolher por padrão de acesso, RPO/RTO e TCO — não por US$/GB isolado.",
    }
    for old, new in replacements.items():
        replace_in_pptx(prs, old, new)

    replace_exact_paragraph(prs.slides[4], "Infraestrutura de terceiros, compartilhada (multi-tenant), acesso via Internet.", "Infraestrutura operada por provedor/multi-tenant; acesso pode ser privado por VPN, Direct Connect ou endpoint.")
    replace_exact_paragraph(prs.slides[7], "•  Reserved Instances — desconto por 1–3 anos", "•  Reserved Instances — benefício de faturamento; reserva de capacidade só em modalidade zonal")
    replace_exact_paragraph(prs.slides[10], "≈ 3.584 TB  ≈  3.670.016 GB", "3.500 TB = 3.500.000 GB")
    replace_exact_paragraph(prs.slides[10], "(padrão binário: 1 PB = 1.024 TB)", "PB decimal; PiB é outra unidade")
    replace_exact_paragraph(prs.slides[10], "•  Preços de lista, ago/2026", "•  Preços de lista consultados em 05/09/2026")
    replace_exact_paragraph(prs.slides[10], "•  Custo de armazenamento apenas", "•  Storage apenas; requests, retrieval, egress e compute separados")
    replace_exact_paragraph(prs.slides[11], "S3 Standard → Glacier Deep Archive representa economia de armazenamento superior a 95%.", "Classes frias reduzem storage, mas cobram retrieval, metadados e permanência mínima.")
    replace_exact_paragraph(prs.slides[11], "Mas Glacier Deep Archive leva 12–48h para recuperar dados — inviável para consultas ativas.", "Deep Archive: recuperação de 12–48 h; inadequado quando o RTO é menor.")
    replace_exact_paragraph(prs.slides[11], "Na prática: camadas combinadas (lifecycle policies), não uma classe única.", "Escolher por acesso, RPO/RTO e TCO; combinar classes por lifecycle.")

    for slide_no, values in [(11, [1150, 9900, 63000]), (12, [74050, 43750, 12600, 3465])]:
        for shape in prs.slides[slide_no - 1].shapes:
            if getattr(shape, "has_chart", False):
                data = CategoryChartData()
                if slide_no == 11:
                    data.categories = ["0–50 TB", "50–500 TB", ">500 TB"]
                    data.add_series("Custo (US$)", values)
                else:
                    data.categories = ["S3 Standard", "Standard-IA", "Glacier Flexible", "Glacier Deep Archive"]
                    data.add_series("US$/mês", values)
                shape.chart.replace_data(data)

    titles = {shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)}
    if "Objeto, bloco, arquivo, archive e banco" not in titles:
        add_slide(prs, "Objeto, bloco, arquivo, archive e banco",
                  "Semântica",
                  "S3: objetos/chaves e API\nEBS: blocos anexados a EC2\nEFS: arquivos NFS compartilhados\nGlacier: classes frias do S3\nRDS/Aurora: engine gerenciada\nRedshift: warehouse analítico",
                  "Escolha",
                  "OLTP autogerenciado: EBS + HA\nCompartilhamento de arquivos: EFS\nData lake/backup: S3\nArquivo raro: Glacier\nDBaaS: RDS/Aurora\nOLAP recorrente: Redshift/Spectrum",
                  13)
        add_slide(prs, "Arquitetura viável para 3,5 PB",
                  "Data lake como base",
                  "S3 com Parquet/ORC e partições\nGlue Catalog\nAthena para ad hoc\nEMR para lote\nRedshift/Spectrum para recorrência\nLifecycle por temperatura",
                  "Limites e operação",
                  "Não cabe em um único Aurora/EBS\nValidar quotas por versão/região\nDefinir crescimento e nº de objetos\nRPO/RTO e testes de recuperação\n10 Gb/s ideal: 32,4 dias\n100 Gb/s ideal: 3,24 dias",
                  13)
        add_slide(prs, "Processamento pode superar o storage",
                  "Athena a US$5/TB escaneado",
                  "1 varredura de 3.500 TB = US$17.500\n1 por mês = US$17.500/mês\n1 por dia (30) = US$525.000/mês\n10% após poda/compressão = US$1.750/execução",
                  "TCO ainda inclui",
                  "PUT/GET/LIST e transições\nRetrieval e metadados\nPermanência mínima\nKMS, catálogo e observabilidade\nReplicação, suporte e egress\nCrescimento e descontos/cotação",
                  13)

    prs.save(PPTX)


if __name__ == "__main__":
    update_docx()
    update_pptx()
    print("DOCX e PPTX de Cloud atualizados.")
