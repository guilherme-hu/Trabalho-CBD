# -*- coding: utf-8 -*-
"""Verificacao por script (v2, apos a segunda rodada adversarial):
(1) refaz TODA conta publicada; (2) confere que os numeros-chave aparecem
identicos no relatorio (.tex) e nos slides (.pptx); (3) confere a checklist
do enunciado; (4) confere invariantes internos do proprio Post-Mortem."""
import io, sys, zipfile
from pathlib import Path
from xml.etree import ElementTree
BASE = Path(__file__).resolve().parent
falhas = []

def chk(nome, calc, publicado, tol=0.01):
    ok = abs(calc - publicado) / (abs(publicado) or 1) < tol
    print(("  OK  " if ok else "  ERRO") + "  %-46s calc=%-14s pub=%s" % (nome, round(calc,4), publicado))
    if not ok: falhas.append(nome)

print("== 1. Contas publicadas ==")
print("-- Fibre Channel: a prova FISICA de que a FCIA publica full-duplex --")
chk("16GFC: 3.200 MB/s em Gb/s",            3200*8/1000, 25.6)
chk("  ... excede a linha de 14,025 Gb/s?", 25.6/14.025, 1.825, tol=0.01)
chk("128GFC: 24.850 MB/s em Gb/s",          24850*8/1000, 198.8)
chk("  ... excede a linha de 112,2 Gb/s?",  198.8/112.2, 1.772, tol=0.01)
if 25.6 <= 14.025 or 198.8 <= 112.2:
    falhas.append("prova fisica full-duplex")
else:
    print("  OK    prova: vazao publicada > taxa de linha -> so pode ser full-duplex")

print("-- Fibre Channel: a convencao de nomenclatura e onde ela quebra --")
chk("16GFC por direcao (24850 -> nao aplica)", 3200/2, 1600)
chk("64GFC por direcao",                       12800/2, 6400)
chk("128GFC por direcao (Gen 8 serial)",       24850/2, 12425)
chk("razao FCIA/convencao ate 64GFC",          12800/6400, 2.00)
chk("razao FCIA/convencao no 128GFC",          24850/12800, 1.941, tol=0.005)
chk("razao das taxas de linha 56,1/28,9",      56.1/28.9, 1.941, tol=0.005)
chk("eficiencia util 64GFC (51,2/57,8)",       (6400*8/1000)/57.8*100, 88.6, tol=0.01)
chk("eficiencia util 128GFC (99,4/112,2)",     (12425*8/1000)/112.2*100, 88.6, tol=0.01)
print("  NOTA  as duas eficiencias coincidem -> a FCIA escalou o 128GFC do 64GFC pela taxa de linha")

print("-- Derivacoes de engenharia --")
chk("16GFC: 14,025 GBd x 64/66 -> Gb/s", 14.025*64/66, 13.60)
chk("16GFC: -> MB/s por direcao",        14.025*64/66/8*1000, 1700, tol=0.005)
chk("8GFC: 8,5 GBd x 8/10 -> Gb/s",      8.5*0.8, 6.8)
chk("Luz em fibra: 3e8/1,5 (m/s)",       3e8/1.5, 2e8)
chk("RTT 100 km em fibra (ms)",          2*100_000/2e8*1000, 1.0)
chk("Propagacao 10 km, um sentido (us)", 10*5, 50)
chk("BB_Credit 10 km @32GFC: KB em voo", 50e-6*3.2e9/1000, 160)
chk("BB_Credit 10 km: quadros de 2.112 B", 160e3/2112, 76, tol=0.02)
chk("Latencia FC 2 saltos (us)",         2*460/1000, 0.92)
chk("FC como % de acesso NVMe de 20us",  0.92/20*100, 4.6)

print("-- AST: os dois fatores, conforme o SGBD --")
chk("256 MiB / 16 KiB (InnoDB)",         256*1024/16, 16384)
chk("256 MiB / 8 KiB (PostgreSQL/Oracle)", 256*1024/8, 32768)
chk("256 MB decimais / 16 KiB",          256*1000*1000/(16*1024), 15625, tol=0.001)

print("-- RAID: penalidade de escrita --")
chk("RAID 1: operacoes por escrita",     2, 2)
chk("RAID 5: 2 leituras + 2 escritas",   2+2, 4)
chk("RAID 6: 3 leituras + 3 escritas",   3+3, 6)

print("-- Fita: LTO-6 do livro x LTO-10 vigente --")
chk("LTO-10: 40 TB a 400 MB/s (h)",      40e12/400e6/3600, 27.8, tol=0.005)
chk("LTO-6: 2,5 TB a 160 MB/s (h)",      2.5e12/160e6/3600, 4.34, tol=0.01)
chk("crescimento de capacidade",         40/2.5, 16)
chk("crescimento de taxa",               400/160, 2.5)
chk("tempo de leitura piorou",           (40e12/400e6)/(2.5e12/160e6), 6.4, tol=0.01)
chk("LTO-10 30TB comprimido 2,5:1",      30*2.5, 75)
chk("LTO-10 40TB comprimido 2,5:1",      40*2.5, 100)
print("  NOTA  LTO-10: 400 x 2,5 = 1000 MB/s, mas o LTO publica 1200 (exigiria 3:1)")
print("        -> divergencia na fonte primaria, documentada na Secao 9.2 do relatorio.")

print("-- Objeto e nuvem --")
print("  NOTA  11 noves e um objetivo de projeto; nao e probabilidade empirica")
chk("Deep Archive: US$/GB-mes -> TB-mes",0.00099*1000, 0.99, tol=0.02)

print("-- Aurora (totais de 30 min, SysBench write-only) --")
chk("Aurora: transacoes (x)",            27_378_000/780_000, 35.1, tol=0.005)
chk("Aurora: IOs por transacao (x)",     7.4/0.95, 7.79, tol=0.005)

print("-- Livros --")
chk("SATA-3: 6 Gb/s com 8b/10b -> MB/s", 6e9*0.8/8/1e6, 600)

print("-- Post-Mortem: invariantes do proprio processo --")
chk("Rodada 1: 5 nao-confirmadas + 11 ressalvas", 5+11, 16)
chk("Rodada 1 + rodada 2 = total de correcoes",   18+15, 33)
chk("Rodada 1: 43+3+1+1 afirmacoes",              43+3+1+1, 48)

# ---------- (2) cruzamento .tex x .pptx ----------
print("\n== 2. Cruzamento relatorio (.tex) x slides (.pptx) ==")
tex = io.open(BASE / "relatorio.tex", encoding="utf-8").read()
try:
    from pptx import Presentation
    prs = Presentation(BASE / "Slides_NAS_SAN_Armazenamento_SBD.pptx")
    partes = []
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame: partes.append(sh.text_frame.text)
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells: partes.append(c.text)
        try: partes.append(sl.notes_slide.notes_text_frame.text)
        except Exception: pass
    ppt = "\n".join(partes)
except Exception as e:
    # Fallback sem dependências: extrai todos os nós <a:t> dos slides e notas.
    # Falhar silenciosamente aqui invalidaria toda a verificação cruzada.
    print("  (python-pptx indisponivel: %s; usando parser OOXML)" % e)
    try:
        partes = []
        with zipfile.ZipFile(BASE / "Slides_NAS_SAN_Armazenamento_SBD.pptx") as pacote:
            nomes = sorted(
                n for n in pacote.namelist()
                if (n.startswith("ppt/slides/slide") or
                    n.startswith("ppt/notesSlides/notesSlide")) and n.endswith(".xml")
            )
            for nome in nomes:
                raiz = ElementTree.fromstring(pacote.read(nome))
                partes.extend(no.text or "" for no in raiz.iter()
                              if no.tag.endswith("}t"))
        if not partes:
            raise ValueError("nenhum texto encontrado no PPTX")
        ppt = "\n".join(partes)
    except Exception as fallback:
        print("  ERRO  nao foi possivel ler o PPTX: %s" % fallback)
        falhas.append("leitura do PPTX")
        ppt = ""

def norm(t):
    # normaliza notacao LaTeX: 27{,}8 -> 27,8 ; 1.700 -> 1700 ; ~ -> espaco
    return t.replace("{,}", ",").replace("{.}", ".").replace("~", " ").replace(".", "")

CHAVES = [
  ("460 ns",             ["460"]),
  ("2000 buffers",       ["2000", "2.000"]),
  ("256 MiB (FAST VP)",  ["256 M"]),
  ("16.384x",            ["16384", "16.384"]),
  ("32.768x",            ["32768", "32.768"]),
  ("128GFC serial",      ["24850", "24.850"]),
  ("128GFC por direcao", ["12425", "12.425"]),
  ("128GFC ISL",         ["25600", "25.600"]),
  ("FC-PI-8 legado",     ["12800", "12.800"]),
  ("linha 56,1 GBd",     ["56,1"]),
  ("RFC 7143",           ["7143"]),
  ("RFC 8881",           ["8881"]),
  ("RFC 3821",           ["3821"]),
  ("INCITS 462-2010",    ["462-2010"]),
  ("EtherType 0x8906",   ["0x8906"]),
  ("IEEE 802.1Qbb",      ["8021Qbb", "802.1Qbb"]),
  ("porta 3260",         ["3260"]),
  ("porta 2049",         ["2049"]),
  ("porta 445",          ["445"]),
  ("porta 548",          ["548"]),
  ("11 noves",           ["99,999999999", "99.999999999", "11 noves"]),
  ("LTO-10 400 MB/s",    ["400 MB/s", "400 MBps"]),
  ("LTO-6 do livro",     ["2,5 TB"]),
  ("Aurora 27.378.000",  ["27378000", "27.378.000"]),
  ("Aurora 30 minutos",  ["30 minutos", "30 min"]),
  ("Aurora SysBench",    ["SysBench"]),
  ("1 ms por 100 km",    ["100 km"]),
  ("27,8 h LTO",         ["27,8"]),
  ("macOS 15.5",         ["15.5"]),
  ("macOS 27",           ["macOS 27"]),
  ("NDMP",               ["NDMP"]),
  ("FCP_XFER_RDY",       ["XFER\\_RDY", "XFER_RDY"]),
  ("ALUA",               ["ALUA"]),
  ("hard vs soft (NFS)", ["hard"]),
  ("33 correcoes",       ["33 corre"]),
]
for nome, alts in CHAVES:
    t_ok = any(a in tex or a in norm(tex) for a in alts)
    p_ok = any(a in ppt or a in norm(ppt) for a in alts)
    status = "  OK  " if (t_ok and p_ok) else "  ERRO"
    if not (t_ok and p_ok): falhas.append("cruzamento: "+nome)
    print("%s  %-22s tex=%s pptx=%s" % (status, nome, t_ok, p_ok))

# ---------- (3) checklist do enunciado ----------
print("\n== 3. Checklist do enunciado presente no relatorio ==")
ITENS = ["SMB/CIFS","NFS","AFP","iSCSI","FCIP","FCoE","FC \\textit{Switch}",
         "Automated Storage Tiering","Object-Based Storage","secundário","terciário",
         "Post-Mortem","DRE","Como funciona"]
for i in ITENS:
    ok = i in tex
    if not ok: falhas.append("checklist: "+i)
    print(("  OK  " if ok else "  ERRO") + "  " + i)

# ---------- (4) higiene do LaTeX ----------
print("\n== 4. Higiene ==")
import re
citadas = set(re.findall(r"\\cite\{([^}]*)\}", tex))
citadas = {c.strip() for grp in citadas for c in grp.split(",")}
declaradas = set(re.findall(r"\\bibitem\{([^}]*)\}", tex))
orfas = declaradas - citadas
faltando = citadas - declaradas
print(("  OK  " if not faltando else "  ERRO") + "  toda \\cite tem \\bibitem  %s" % (sorted(faltando) or ""))
print(("  OK  " if not orfas else "  AVISO") + "  referencias declaradas e nao citadas: %s" % (sorted(orfas) or "nenhuma"))
if faltando: falhas.append("bibitem faltando")

print("\n== RESUMO ==")
if falhas:
    print("FALHAS (%d): %s" % (len(falhas), falhas)); sys.exit(1)
print("Tudo confere.")
