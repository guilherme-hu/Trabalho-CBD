#!/usr/bin/env python3
"""Gera os PPTX finais de DAS/RAID e Cloud/AWS com python-pptx."""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
W, H = Inches(13.333), Inches(7.5)
C = {"dark":"14181F","dark2":"232A35","ink":"14181F","muted":"55606E",
     "white":"FFFFFF","soft":"F4F5F7","orange":"EB6834","blue":"2A78D6",
     "line":"DFE3E8","green":"178B67"}
NAMES = [
    ("Bernardo Brandão Pozzato Carvalho Costa","123289593"),
    ("Enzo de Carvalho Sampaio","123386206"),
    ("Gabriel Schmitz Corrêa Rizawinsk","123225573"),
    ("Guilherme En Shih Hu","123224674"),
    ("Raphael Henrique da Silva Pereira","123311073"),
    ("Vivian Maria da Silva e Souza","123205793"),
]

def rgb(value): return RGBColor.from_string(value)

def textbox(slide, text, x, y, w, h, size=15, color="ink", bold=False,
            face="Calibri", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear(); box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = valign
    box.text_frame.margin_left = box.text_frame.margin_right = Inches(.04)
    p = box.text_frame.paragraphs[0]; p.text = text; p.alignment = align
    p.font.name = face; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = rgb(C.get(color,color))
    return box

def new_prs(title):
    prs=Presentation(); prs.slide_width=W; prs.slide_height=H
    prs.core_properties.title=title; prs.core_properties.author="Grupo CBD — UFRJ"
    return prs

def blank(prs, dark=False):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(C["dark"] if dark else C["white"])
    return s

def heading(s, title, subtitle, page, footer):
    textbox(s,title,.72,.42,11.9,.55,27,"ink",True,"Cambria")
    if subtitle: textbox(s,subtitle,.72,1.03,11.9,.35,13,"muted")
    line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.72), Inches(6.87), Inches(11.9), Pt(.7))
    line.fill.solid(); line.fill.fore_color.rgb=rgb(C["line"]); line.line.fill.background()
    textbox(s,footer,.72,6.96,9,.23,9,"muted")
    textbox(s,str(page),12,6.96,.6,.23,9,"muted",align=PP_ALIGN.RIGHT)

def slide(prs,title,subtitle,page,footer):
    s=blank(prs); heading(s,title,subtitle,page,footer); return s

def card(s,x,y,w,h,title,body,color="blue",fs=12):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(C["soft"]); sh.line.color.rgb=rgb(C[color]); sh.line.width=Pt(1)
    textbox(s,title,x+.18,y+.11,w-.36,.32,14,color,True)
    textbox(s,body,x+.18,y+.48,w-.36,h-.57,fs,"ink")

def bullet_list(s, items, y=1.55, size=15):
    box=s.shapes.add_textbox(Inches(.78),Inches(y),Inches(11.75),Inches(5.05))
    tf=box.text_frame; tf.clear(); tf.word_wrap=True
    for i,item in enumerate(items):
        para=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.text=item; para.level=0; para.font.name="Calibri"; para.font.size=Pt(size)
        para.font.color.rgb=rgb(C["ink"]); para.space_after=Pt(9); para.text="•  "+para.text

def table(s, headers, rows, y, widths, size=11):
    total_rows=len(rows)+1; shape=s.shapes.add_table(total_rows,len(headers),Inches(.72),Inches(y),Inches(11.9),Inches(min(4.95,.42*total_rows)))
    tbl=shape.table
    for j,width in enumerate(widths): tbl.columns[j].width=Inches(width)
    for i,row in enumerate([headers]+rows):
        for j,value in enumerate(row):
            cell=tbl.cell(i,j); cell.text=str(value); cell.margin_left=cell.margin_right=Inches(.06); cell.margin_top=cell.margin_bottom=Inches(.04)
            cell.fill.solid(); cell.fill.fore_color.rgb=rgb(C["dark2"] if i==0 else C["white"])
            cell.border if False else None
            for para in cell.text_frame.paragraphs:
                para.font.name="Calibri"; para.font.size=Pt(size); para.font.bold=(i==0)
                para.font.color.rgb=rgb(C["white"] if i==0 else C["ink"])
                para.alignment=PP_ALIGN.LEFT
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    return tbl

def cover(prs, title, subtitle, group):
    s=blank(prs,True)
    textbox(s,"CONSTRUÇÃO DE BANCO DE DADOS · UFRJ",.72,.68,11.8,.3,11,"orange",True)
    textbox(s,title,.72,1.2,7.8,1.5,41,"white",True,"Cambria")
    textbox(s,subtitle,.72,2.95,8.7,.6,16,"CBD2DC")
    table_data=[("Nome completo","DRE")]+NAMES
    shape=s.shapes.add_table(len(table_data),2,Inches(.72),Inches(3.72),Inches(8.9),Inches(2.58))
    tbl=shape.table; tbl.columns[0].width=Inches(6.8); tbl.columns[1].width=Inches(2.1)
    for i,row in enumerate(table_data):
        for j,val in enumerate(row):
            c=tbl.cell(i,j); c.text=val; c.fill.solid(); c.fill.fore_color.rgb=rgb(C["dark2"] if i==0 else C["dark"])
            for q in c.text_frame.paragraphs:
                q.font.name="Calibri"; q.font.size=Pt(11); q.font.bold=(i==0); q.font.color.rgb=rgb(C["white"])
    textbox(s,group,9.85,4.0,2.65,1.0,14,"white",True,align=PP_ALIGN.RIGHT)
    return s

def three_cards(s, cards, y=1.62, h=3.35):
    for i,(head,body,color) in enumerate(cards): card(s,.72+i*4.13,y,3.65,h,head,body,color)

def make_das():
    prs=new_prs("DAS e RAID para SBD"); footer="DAS e RAID para SBD — UFRJ"; n=0
    cover(prs,"DAS e\narquiteturas RAID","Direct Attached Storage · Interfaces, desempenho, confiabilidade e decisão","GRUPO 3\nSetembro de 2026")
    def S(t,sub=""):
        nonlocal n; n+=1; return slide(prs,t,sub,n,footer)
    s=S("DAS e RAID são duas decisões diferentes","DAS conecta; RAID distribui dados")
    card(s,.72,1.65,5.75,2.25,"DAS","Conexão direta ao host, sem rede de armazenamento. O host vê blocos. Simples, porém limitado por baias, portas e domínio do servidor.")
    card(s,6.85,1.65,5.75,2.25,"RAID","Layout lógico entre discos para paralelismo, redundância ou ambos. Pode existir em software, controladora ou array.")
    card(s,.72,4.25,11.88,1.25,"Decisão do SGBD","Medir latência de flush, filas, cache, falhas, rebuild, RPO e RTO — não apenas Gb/s.","orange")
    s=S("Interfaces DAS","Taxa nominal não é vazão sustentada")
    table(s,["Interface","Nominal","Recursos","Uso/estado"],[["SATA 6 Gb/s","6 Gb/s","ponto a ponto; NCQ","custo/capacidade"],["eSATA","3/6 Gb/s","externo; sem energia original","legado"],["SAS-4","22,5 Gb/s/lane","dual-port; expansores","empresarial"],["FireWire 800","800 Mb/s","peer-to-peer","legado"],["USB 3.2","até 20 Gb/s","host; UASP","backup/dev"],["USB4 v2","até 80 Gb/s","túneis; negociação","SSD externo"],["HBA","depende","adaptador SAS/FC etc.","não é protocolo"]],1.5,[2.05,2.05,4.05,3.75],11)
    textbox(s,"USB4 120 Gb/s é opcional/assimétrico. SATA-IO recomenda “SATA 6 Gb/s”, não “SATA III”.",.72,5.95,11.9,.38,11.5,"orange",True,align=PP_ALIGN.CENTER)
    s=S("Paralelismo × redundância","Os mesmos discos financiam os dois objetivos")
    card(s,.72,1.65,5.75,2.35,"Paralelismo","Striping pode elevar throughput de uma requisição ou atender I/Os independentes. Ganho N× é um limite ideal.")
    card(s,6.85,1.65,5.75,2.35,"Redundância","Espelho ou paridade mantém informação extra. Disponibilidade é continuar; durabilidade é preservar.")
    card(s,.72,4.4,11.88,1.2,"Pergunta correta","Qual falha é tolerada, quanto dura o rebuild e qual desempenho resta no estado degradado?","orange")
    s=S("Data striping","A granularidade muda o paralelismo")
    three_cards(s,[("Bit-level","Bits espalhados; discos em lockstep. RAID 2 acadêmico.","blue"),("Byte/fino","RAID 3: todos os discos participam da I/O; streaming.","blue"),("Block-level","Chunks inteiros; I/Os independentes. RAID 0/4/5/6.","green")],1.65,3.1)
    textbox(s,"Stripe unit = chunk por disco · stripe = conjunto · stripe width = dados antes da paridade",.72,5.35,11.9,.4,14,"blue",True,align=PP_ALIGN.CENTER)
    s=S("Mirroring, shadowing e paridade","Termos parecidos em camadas diferentes")
    three_cards(s,[("Mirroring","Cópia completa; leituras distribuíveis; toda escrita vai aos membros.","blue"),("Shadowing","Pode nomear espelho histórico. Shadow paging e snapshot não são RAID 1.","orange"),("Paridade","P = D₀ ⊕ D₁ ⊕ …; RAID 6 adiciona equação independente.","green")],1.65,3.35)
    s=S("Write penalty e write hole","A pequena escrita de paridade custa caro")
    card(s,.72,1.62,5.75,3.2,"RAID 5: read-modify-write","Ler dado antigo + ler paridade antiga + gravar dado novo + gravar paridade nova = 4 I/Os de membro. RAID 6 pode chegar a 6.")
    card(s,6.85,1.62,5.75,3.2,"Consistência","Stripe completo evita leituras antigas. Queda entre dado e paridade causa write hole; cache protegido, journal/bitmap ou RAID-Z tratam o risco.","orange")
    s=S("MTTF, MTBF, MTTR e MTTDL","Modelo útil, hipóteses explícitas")
    table(s,["Métrica","Pergunta"],[["MTTF","tempo médio até falha"],["MTTR","detecção + troca + reconstrução"],["MTBF","MTTF + MTTR em convenção didática"],["MTTDL","tempo médio até perda de dados"],["Disponibilidade","MTTF/(MTTF+MTTR)"]],1.52,[2.8,9.1],13)
    textbox(s,"RAID 5 idealizado: MTTDL ≈ MTTF² / [N(N−1)MTTR]",.72,4.5,11.9,.45,22,"blue",True,"Cambria",PP_ALIGN.CENTER)
    textbox(s,"Não inclui URE, falha correlacionada, setor latente nem carga de rebuild.",.72,5.28,11.9,.4,14,"orange",True,align=PP_ALIGN.CENTER)
    s=S("RAID 0, 1, 5 e 6","Capacidade, falha e escrita")
    table(s,["Nível","Útil","Tolera","Escrita pequena","Uso"],[["0","N·S","0","sem paridade","scratch recriável"],["1","S/par","1/par","duas cópias","logs/sistema"],["5","(N−1)·S","1","até 4 I/Os","leitura; rebuild analisado"],["6","(N−2)·S","2","até 6 I/Os","grandes conjuntos/leitura"]],1.58,[1.5,2.1,1.45,2.65,4.2],12.5)
    card(s,.72,4.55,11.88,1.22,"RAID 5 não é “uso geral” automaticamente","Discos grandes ampliam reconstrução e exposição a URE; medir no hardware real.","orange")
    s=S("RAID 2, 3 e 4","Referência histórica")
    three_cards(s,[("RAID 2","Bits + Hamming dedicado; ECC já existe no drive.","blue"),("RAID 3","Striping fino + paridade dedicada; lockstep.","blue"),("RAID 4","Blocos + paridade dedicada; gargalo de escrita.","blue")],1.65,3.25)
    textbox(s,"RAID 5 distribui a paridade e remove o ponto quente fixo do RAID 4.",.72,5.35,11.9,.42,15,"green",True,align=PP_ALIGN.CENTER)
    s=S("RAID 0+1 × RAID 10","A ordem muda o domínio de falha")
    card(s,.72,1.65,5.75,3.1,"0+1: espelho de stripes","Após uma falha, um stripe inteiro pode ficar degradado. Segunda falha no stripe sobrevivente causa perda; rebuild pode copiar o conjunto.")
    card(s,6.85,1.65,5.75,3.1,"10: stripe de espelhos","Só o par afetado fica exposto. Falhas em pares distintos podem coexistir; rebuild é localizado.","green")
    textbox(s,"Mesma capacidade de 50%; RAID 10 costuma ter comportamento de falha/rebuild melhor.",.72,5.25,11.9,.4,15,"orange",True,align=PP_ALIGN.CENTER)
    s=S("Níveis não padrão","Rótulos de produto e composições")
    table(s,["Nome","Origem/ideia","Leitura"],[["RAID 1.5","HighPoint","marca; RAID 1 pode balancear leitura"],["RAID 7","Storage Computer Corp.","produto, não sucessor"],["RAID-DP","NetApp/WAFL","dupla paridade integrada"],["RAID-S","EMC Symmetrix","depende da geração"],["Matrix RAID","Intel","volumes compartilham discos"],["RAID-Z","ZFS","paridade + checksums; evita write hole"]],1.52,[2.1,4.1,5.7],11.5)
    s=S("Recomendação por componente","Workload e RPO/RTO decidem")
    table(s,["Componente","Candidato","Critério"],[["WAL/redo/undo","RAID 1/10","flush e cache protegido"],["Dados OLTP","RAID 10; 6 se leitura","IOPS, tail latency, rebuild"],["DW/OLAP","RAID 6/60 ou distribuído","throughput e rebuild em carga"],["Backup local","RAID 6/Z2 no repositório","cópia independente permanece"],["Temp/scratch","RAID 0 condicionado","recriável dentro do RTO"],["DAS externo","USB4/SAS externo","energia, flush, recuperação"]],1.52,[3.05,3.35,5.5],11.8)
    s=S("RAID não é backup","Controles complementares")
    card(s,.72,1.62,5.75,3.35,"RAID: continuidade","Tolera falhas previstas\nMantém serviço durante rebuild\nNão preserva versão anterior\nNão cobre exclusão/ransomware")
    card(s,6.85,1.62,5.75,3.35,"Backup: recuperação","Cópia independente/versionada\nPreferência off-site/imutável\nRPO e RTO definidos\nRestauração testada","green")
    textbox(s,"Hot spare reduz a janela de exposição; ele também não é backup.",.72,5.45,11.9,.4,15,"orange",True,align=PP_ALIGN.CENTER)
    s=S("Conclusões","O rótulo é só o começo")
    bullet_list(s,["DAS conecta o host; RAID distribui os dados","SATA prioriza custo; SAS acrescenta caminhos; HBA é adaptador","Block striping permite I/O independente; bit/fino opera em lockstep","RAID 10 compra escrita/rebuild com 50%; RAID 6 compra tolerância com paridade","MTTDL orienta; URE, correlação e rebuild exigem teste","RAID, backup e recuperação de desastre são complementares"],1.52,14.5)
    s=S("Post-Mortem","Uso crítico e rastreável de IA")
    three_cards(s,[("Prompts","Enunciado, geração, pesquisa e revisão final registrados.","blue"),("11 correções","Estrutura, interfaces, métricas, níveis, backup e autoria.","orange"),("Pendência humana","Confirmar grupo e registrar divisão real; nada foi inventado.","green")],1.65,3.5)
    end=blank(prs,True); textbox(end,"Perguntas e debate",.72,1.45,11.9,.8,40,"white",True,"Cambria",PP_ALIGN.CENTER); textbox(end,"SSD grande torna RAID 5 mais seguro — ou amplia a janela de risco?\nQuando a replicação do SGBD deve substituir RAID?",1.25,2.7,10.8,1.5,20,"D6DCE5",align=PP_ALIGN.CENTER); textbox(end,"Grupo 3 · Setembro de 2026",.72,6.55,11.9,.3,11,"orange",align=PP_ALIGN.CENTER)
    prs.save(ROOT/"DAS e RAID"/"DAS_RAID_Slides.pptx")

def make_cloud():
    prs=new_prs("Cloud Computing e Cloud Storage — AWS");footer="Cloud Computing e Cloud Storage — AWS";n=0
    cover(prs,"Cloud Computing\ne Cloud Storage","Estudo de caso AWS · arquitetura e custo para 3,5 PB","GRUPO 3\nSetembro de 2026")
    def S(t,sub=""):
        nonlocal n;n+=1;return slide(prs,t,sub,n,footer)
    s=S("O que torna uma infraestrutura “nuvem”?","Cinco características do NIST");bullet_list(s,["Autoatendimento sob demanda","Amplo acesso por rede","Pool de recursos compartilhado","Elasticidade rápida","Serviço mensurável (pay-as-you-go)"],1.6,18);textbox(s,"Virtualização ajuda, mas não satisfaz sozinha a definição.",.72,5.75,11.9,.4,14,"orange",True,align=PP_ALIGN.CENTER)
    s=S("Cloud Storage tem três semânticas","Objeto, bloco e arquivo não são intercambiáveis");three_cards(s,[("Objeto — S3","Chave + metadados + API. Data lake, conteúdo e backup.","blue"),("Bloco — EBS","Volume para EC2; o host monta o sistema de arquivos.","green"),("Arquivo — EFS","Diretórios NFS compartilháveis; modos e classes próprios.","orange")],1.65,3.4);textbox(s,"Glacier são classes frias do S3 — não um disco montável.",.72,5.45,11.9,.4,15,"orange",True,align=PP_ALIGN.CENTER)
    s=S("Modelos de serviço","A fronteira muda; a responsabilidade pelos dados permanece");table(s,["Modelo","Exemplo","Provedor opera","Cliente decide"],[["IaaS","EC2 + EBS","infra/hipervisor","SO, SGBD, HA, backup"],["PaaS","Elastic Beanstalk","infra/plataforma","código, dados, configuração"],["SaaS","QuickSight","aplicação/pilha","conteúdo, usuários, acesso"],["DBaaS","RDS/Aurora","engine/tarefas","esquema, SQL, IAM"],["DWaaS","Redshift","MPP/infra","modelo, cargas, WLM"]],1.52,[1.45,2.25,3.8,4.4],11.5)
    s=S("Pública, privada e híbrida","Implantação, não endereço IP");three_cards(s,[("Pública","Provedor multi-tenant com isolamento; pode usar conexão privada.","blue"),("Privada","Dedicada à organização; controle e operação/capacidade próprios.","orange"),("Híbrida","Integra rede, identidade, gestão e dados entre ambientes.","green")],1.65,3.55)
    s=S("Sete critérios de avaliação","Cada porcentagem precisa de contexto");table(s,["Critério","Pergunta"],[["Acesso","Internet, VPN, Direct Connect ou endpoint?"],["Disponibilidade","qual SLA e desenho por zona/região?"],["Latência","distância + serviço + fila + aplicação?"],["Escala","limite por recurso e elasticidade?"],["Durabilidade","o que é preservado; quais ameaças ficam fora?"],["Segurança","quem configura IAM, KMS, rede e logs?"],["Custo","capacidade + operações + compute + egress + suporte?"]],1.46,[2.25,9.65],11.8)
    s=S("Durabilidade ≠ disponibilidade ≠ backup","S3 é projetado para 11 noves de durabilidade");three_cards(s,[("Durabilidade","Probabilidade de preservar objetos; múltiplas zonas.","blue"),("Disponibilidade","Probabilidade de responder; SLA depende da classe.","orange"),("Recuperação","Versionamento, Object Lock, cópia e teste de restore.","green")],1.65,3.4);textbox(s,"11 noves não prometem zero downtime nem retorno a uma versão anterior.",.72,5.45,11.9,.4,15,"orange",True,align=PP_ALIGN.CENTER)
    s=S("AWS: computação e armazenamento","Escolha pela semântica");table(s,["Serviço","Função","Cuidado"],[["EC2","máquina virtual/IaaS","cliente opera SO/SGBD/HA"],["S3","objeto/data lake","requests, classes, egress, formato"],["EBS","bloco para EC2","16 TiB por gp3; IOPS/throughput"],["EFS","arquivo NFS","modo/classe e acesso"],["Glacier","objeto frio","duração, metadados e restauração"]],1.52,[2,4.15,5.75],12)
    s=S("AWS: bancos e data warehouse","Gerenciado não significa sem decisões");table(s,["Serviço","Modelo","Adequação","Cuidado"],[["RDS","DBaaS","relacional multi-engine","HA depende de Multi-AZ"],["Aurora","DBaaS","MySQL/PostgreSQL compatível","I/O, capacidade, limites"],["DynamoDB","NoSQL","chave–valor/documento","partição e consistência"],["Redshift","DWaaS","OLAP colunar/MPP","compute, RMS, WLM"]],1.58,[2.1,2.25,3.5,4.05],11.5);textbox(s,"3,5 PB exigem arquitetura distribuída; não cabem em uma única instância ou volume.",.72,4.65,11.9,.5,15,"orange",True,align=PP_ALIGN.CENTER)
    s=S("3,5 PB: unidade e premissas","Preço sem região, data e workload é falsa precisão");card(s,.72,1.62,5.75,2.2,"Conversão","3,5 PB decimais = 3.500 TB = 3.500.000 GB. Não misturar PB com PiB/GiB.");card(s,6.85,1.62,5.75,2.2,"Escopo","us-east-1 · 05/09/2026 · lista pública · sem impostos, suporte ou desconto.");card(s,.72,4.2,11.88,1.35,"Variáveis","Crescimento, compressão, objetos, temperatura, consultas, região, RPO/RTO, retenção e egress.","orange")
    s=S("Quanto custa armazenar 3,5 PB?","Capacidade mensal; serviços não equivalentes");table(s,["Classe/serviço","Preço","US$/mês","Condição"],[["S3 Standard","faixas 0,023/0,022/0,021","74.050","frequente"],["Standard-IA","0,0125/GB","43.750","30 d; retrieval; 128 KB"],["Glacier Flexible","0,0036/GB","12.600","90 d; restauração"],["Deep Archive","0,00099/GB","3.465","180 d; 12–48 h"],["Redshift RMS","0,024/GB","84.000","compute separado"],["EBS gp3","0,08/GB","280.000","199+ volumes; IOPS"]],1.48,[2.7,2.7,2.05,4.45],11.3);textbox(s,"S3 Standard = 50k×0,023 + 450k×0,022 + 3M×0,021 = US$ 74.050/mês",.72,5.92,11.9,.35,12.2,"blue",True,align=PP_ALIGN.CENTER)
    s=S("Processar pode custar mais que guardar","Athena a US$ 5/TB escaneado");three_cards(s,[("1 scan integral","3.500 TB × US$ 5\n= US$ 17.500","blue"),("1 scan/dia","US$ 17.500 × 30\n= US$ 525.000/mês","orange"),("10% lido","Parquet + partição\n= US$ 1.750/consulta","green")],1.65,3.35);textbox(s,"Formato físico, compressão e partição são decisões financeiras.",.72,5.45,11.9,.4,16,"blue",True,align=PP_ALIGN.CENTER)
    s=S("Mover 3,5 PB também custa tempo","Limite físico antes de overhead");textbox(s,"t = volume × 8 / banda",.72,1.65,11.9,.5,26,"blue",True,"Cambria",PP_ALIGN.CENTER);card(s,1.45,2.65,4.8,2.05,"10 Gb/s contínuos","≈ 32,4 dias");card(s,7.05,2.65,4.8,2.05,"100 Gb/s contínuos","≈ 3,24 dias","green");textbox(s,"Protocolo, criptografia, contenção e retransmissão aumentam o tempo.",.72,5.2,11.9,.5,14,"orange",True,align=PP_ALIGN.CENTER)
    s=S("Arquitetura recomendada","S3 durável; compute escolhido pela consulta");table(s,["Camada","Desenho","Motivo"],[["Dados","S3 Parquet/ORC particionado","reduz scan; desacopla compute"],["Catálogo","Glue","esquema e partições"],["Ad hoc","Athena","serverless por bytes lidos"],["Lote","EMR/Spark","processamento distribuído"],["Recorrente","Redshift + Spectrum","MPP, cache, concorrência"],["Lifecycle","Standard → IA → Glacier","custo segue temperatura"],["Operacional","RDS/Aurora/DynamoDB separado","transações em subconjunto"]],1.45,[2.25,4.65,5],11.3)
    s=S("Três cenários para os mesmos 3,5 PB","A pergunta de custo precisa de um SLO");three_cards(s,[("Arquivo","Deep Archive + Object Lock + cópia; RTO 12–48 h.","blue"),("Data lake","S3 Parquet + Glue + Athena/EMR; consultas seletivas.","green"),("Analytics recorrente","S3 + Redshift; baixa latência e concorrência.","orange")],1.65,3.65)
    s=S("Conclusões","O dado cabe; a arquitetura precisa fazer sentido");bullet_list(s,["Cloud exige autoatendimento, pool, elasticidade e mensuração","Objeto, bloco e arquivo oferecem semânticas diferentes","S3 é base plausível para analytics; bancos operacionais recebem subconjuntos","Capacidade S3: US$ 3.465–74.050/mês nas premissas","Scan diário integral no Athena: US$ 525.000/mês","Decisão exige workload, RPO/RTO, segurança, recuperação, egress e TCO"],1.52,14.5)
    s=S("Post-Mortem","IA como rascunho e objeto de auditoria");three_cards(s,[("Prompts","Enunciado, AWS, custo e revisão final registrados.","blue"),("10 correções","Ordem, repetição, taxonomia, custo, limites e autoria.","orange"),("Pendência","Confirmar Grupo 3/7 e divisão humana real.","green")],1.65,3.5)
    end=blank(prs,True);textbox(end,"Perguntas e debate",.72,1.45,11.9,.8,40,"white",True,"Cambria",PP_ALIGN.CENTER);textbox(end,"Que fração dos 3,5 PB precisa de milissegundos?\nQuando um warehouse fica mais barato que consultas serverless?",1.25,2.7,10.8,1.5,20,"D6DCE5",align=PP_ALIGN.CENTER);textbox(end,"Grupo 3 · Setembro de 2026",.72,6.55,11.9,.3,11,"orange",align=PP_ALIGN.CENTER)
    prs.save(ROOT/"Cloud Computing & Cloud Storages"/"Cloud_Computing_AWS_Slides.pptx")

if __name__ == "__main__":
    make_das(); make_cloud(); print("Slides finais gerados.")
